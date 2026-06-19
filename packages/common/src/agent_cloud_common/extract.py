"""按扩展名把文件抽取成文本,供 read_file 工具使用。

纯文本/代码/csv/json 等直接按 UTF-8 读;PDF/Word/PPT/Excel 这类二进制文档用对应库
抽取成轻结构文本(表格转 '|' 分隔、分页/分 sheet 标注),让模型能直接理解上传的文档
内容,而不必自己写解析代码。解析库按需(lazy)导入。

加固(对抗审查后):
- 输入大小双闸:文件硬上限 + zip 类(docx/pptx/xlsx)解压后总大小上限。文档抽取会把内容
  载入内存(docx/pptx 构造即全解压 XML),一个几 KB 的 zip 炸弹可放大到 GB 级撑爆沙箱。
- 库内部异常只回**文件名**(不含完整 host 路径,避免泄露目录结构),由本模块统一包成
  RuntimeError;调用方(sandbox.run_tool / worker.LocalToolExecutor)再转 is_error 交回模型。
- 抽不出实质文本(扫描件/纯图 PDF、空文档)时返回明确提示而非空串,免得模型误判"文件是空的"。

定位:这是"读懂内容"的轻量抽取,不追求还原排版。已知局限——扫描件需 OCR(不含)、复杂
表格与多栏 PDF 会走样、xlsx 取的是缓存计算值。复杂场景模型可改用 bash 工具(沙箱已装这些库)
自行精解析。
"""

from __future__ import annotations

import zipfile
from pathlib import Path

# 自动抽取文本的二进制文档格式(OOXML / PDF)。其余后缀按纯文本读。
DOC_SUFFIXES = frozenset({".pdf", ".docx", ".pptx", ".xlsx", ".xlsm"})

# 旧二进制 Office 格式(非 OOXML zip),python 库不认,需 LibreOffice 转换,第一版不支持。
_LEGACY_OFFICE = frozenset({".doc", ".xls", ".ppt"})

# 图片格式:read_file 不解析像素;返回提示引导"作为附件上传给 vision 模型"(spec:
# image-understanding),而非让模型以为文件损坏、反复 read_file 浪费回合。
IMAGE_SUFFIXES = frozenset({".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"})
_IMAGE_HINT = (
    "is an image; read_file does not parse image pixels. To let the model see this image, "
    "attach it in the chat instead (requires a vision-capable model)."
)

# 抽取的内存防线:输入文件硬上限;zip 类额外限制解压后总大小(防 zip 炸弹高压缩比放大)。
_MAX_INPUT_BYTES = 25 * 1024 * 1024
_MAX_UNCOMPRESSED_BYTES = 100 * 1024 * 1024

_BLANK_HINT = (
    "no extractable text — it may be a scanned/image-only PDF (needs OCR) or an empty document"
)


def extract_text(path: Path) -> str:
    """读取文件文本。文档类(pdf/docx/pptx/xlsx)自动抽取,其余按 UTF-8 文本读。"""
    suffix = path.suffix.lower()
    if suffix in _LEGACY_OFFICE:
        raise RuntimeError(
            f"{path.name}: legacy binary Office format ({suffix}) is not supported; "
            "convert it to .docx/.xlsx/.pptx and upload again."
        )
    if suffix in IMAGE_SUFFIXES:
        # 图片只回引导提示、不读像素,故置于大小检查前(不载入文件、不受 25MB 限)。
        return f"{path.name} {_IMAGE_HINT}"

    _guard_input_size(path)

    if suffix not in DOC_SUFFIXES:
        # 纯文本/代码/csv/json/md/...:按 UTF-8 读。二进制 → 友好提示而非裸 UnicodeDecodeError。
        try:
            return path.read_text()
        except UnicodeDecodeError as exc:
            raise RuntimeError(
                f"{path.name}: not UTF-8 text and not a supported document type "
                "(pdf/docx/pptx/xlsx); it looks binary. If it is a document, convert it to a "
                "supported format, or parse it yourself with the bash tool."
            ) from exc

    if suffix != ".pdf":  # docx/pptx/xlsx 都是 zip
        _guard_uncompressed_size(path)

    try:
        if suffix == ".pdf":
            text = _from_pdf(path)
        elif suffix == ".docx":
            text = _from_docx(path)
        elif suffix == ".pptx":
            text = _from_pptx(path)
        else:  # .xlsx / .xlsm
            text = _from_xlsx(path)
    except RuntimeError:
        raise  # 已是不含 host 路径的友好错误,直接上抛
    except Exception as exc:
        # 库内部异常的 str 可能含完整 host 路径(如 python-docx 的 PackageNotFoundError);
        # 只回文件名 + 异常类型,避免把 host 目录结构泄露给模型。
        raise RuntimeError(
            f"{path.name}: could not extract text ({type(exc).__name__}); "
            f"the file may be corrupt or not a valid {suffix} document."
        ) from exc

    if _is_blank(text):
        return f"{path.name}: {_BLANK_HINT}."
    return text


def _guard_input_size(path: Path) -> None:
    size = path.stat().st_size
    if size > _MAX_INPUT_BYTES:
        raise RuntimeError(
            f"{path.name}: file is too large to read here "
            f"({size // (1024 * 1024)} MB > {_MAX_INPUT_BYTES // (1024 * 1024)} MB limit); "
            "use the bash tool (e.g. head/grep) to read part of it."
        )


def _guard_uncompressed_size(path: Path) -> None:
    # 小压缩包可解压成 GB 级(zip 炸弹)→ 抽取时 OOM。按 zip header 声明的解压总大小预检
    # (标准炸弹靠高压缩比、诚实声明大 size)。顺带把"不是合法 zip/OOXML"挡成友好错误。
    try:
        with zipfile.ZipFile(path) as zf:
            total = sum(info.file_size for info in zf.infolist())
    except zipfile.BadZipFile as exc:
        raise RuntimeError(
            f"{path.name}: not a valid {path.suffix} file (corrupt or not a real OOXML document)."
        ) from exc
    if total > _MAX_UNCOMPRESSED_BYTES:
        raise RuntimeError(
            f"{path.name}: document expands too large to extract safely "
            f"({total // (1024 * 1024)} MB uncompressed)."
        )


def _is_blank(text: str) -> bool:
    """去掉结构标记行(--- page/slide/sheet ---)后若无实质文本,视为空(扫描件/空文档)。"""
    body = "\n".join(ln for ln in text.splitlines() if not ln.startswith("--- "))
    return not body.strip()


def _from_pdf(path: Path) -> str:
    import pdfplumber

    parts: list[str] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            parts.append(f"--- page {i} ---\n{text}".rstrip())
    return "\n\n".join(parts).strip()


def _from_docx(path: Path) -> str:
    import docx  # python-docx

    document = docx.Document(str(path))
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        parts.append(_table_to_text([[c.text for c in row.cells] for row in table.rows]))
    return "\n".join(parts).strip()


def _from_pptx(path: Path) -> str:
    from pptx import Presentation  # python-pptx

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                parts.append(
                    _table_to_text([[c.text for c in row.cells] for row in shape.table.rows])
                )
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(f"[notes] {notes}")
    return "\n".join(parts).strip()


def _from_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    # data_only:取 Excel 缓存的公式计算值(用户存的文件通常有缓存);read_only 省内存。
    # openpyxl 自生成、从未被 Excel 算过的公式单元格其缓存值为 None → 显示为空。
    wb = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            rows = [
                ["" if v is None else str(v) for v in row]
                for row in ws.iter_rows(values_only=True)
            ]
            rows = [r for r in rows if any(c.strip() for c in r)]
            if rows:
                parts.append(f"--- sheet: {ws.title} ---")
                parts.append("\n".join("\t".join(r) for r in rows))
    finally:
        wb.close()
    return "\n".join(parts).strip()


def _table_to_text(rows: list[list[str]]) -> str:
    """表格转 markdown 风格的 '|' 分隔文本(轻结构,模型易读)。"""
    return "\n".join(
        "| " + " | ".join(c.replace("\n", " ").strip() for c in row) + " |" for row in rows
    )
