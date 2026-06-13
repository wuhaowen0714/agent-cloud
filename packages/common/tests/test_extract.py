"""extract_text:文档类(pdf/docx/pptx/xlsx)自动抽取文本,其余按 UTF-8 纯文本。

含对抗审查后补的回归:zip 炸弹/大小闸、损坏文件不泄露 host 路径、扫描件/空文档提示、
XXE 外部实体不解析。
"""

from __future__ import annotations

import zipfile

import pytest
from agent_cloud_common import extract as extract_mod
from agent_cloud_common.extract import extract_text

# ---- 纯文本路径 ----


def test_plain_text_returned_as_is(tmp_path):
    p = tmp_path / "note.txt"
    p.write_text("hello\nworld")
    assert extract_text(p) == "hello\nworld"


def test_csv_returned_as_is(tmp_path):
    p = tmp_path / "data.csv"
    p.write_text("a,b\n1,2")
    assert extract_text(p) == "a,b\n1,2"


def test_unknown_text_extension_read_as_text(tmp_path):
    p = tmp_path / "app.log"
    p.write_text("log line")
    assert extract_text(p) == "log line"


def test_binary_non_utf8_raises_friendly_error(tmp_path):
    p = tmp_path / "blob.bin"
    p.write_bytes(b"\xff\xfe\x00\x01\x02\x03")
    with pytest.raises(RuntimeError, match="binary"):
        extract_text(p)


def test_legacy_office_raises_friendly_error(tmp_path):
    p = tmp_path / "old.doc"
    p.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1")  # OLE2 magic header
    with pytest.raises(RuntimeError, match=r"\.doc|legacy"):
        extract_text(p)


# ---- 文档抽取主路径 ----


def test_docx_extracts_paragraphs_and_tables(tmp_path):
    import docx

    document = docx.Document()
    document.add_paragraph("Hello from docx")
    table = document.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Name"
    table.rows[0].cells[1].text = "Qty"
    table.rows[1].cells[0].text = "Apple"
    table.rows[1].cells[1].text = "5"
    p = tmp_path / "doc.docx"
    document.save(str(p))

    out = extract_text(p)
    assert "Hello from docx" in out
    assert "Name" in out and "Apple" in out and "5" in out
    assert "|" in out  # 表格转 '|' 分隔的轻结构


def test_xlsx_extracts_values_and_sheet_headers(tmp_path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["name", "qty"])
    ws.append(["apple", 5])
    second = wb.create_sheet("More")
    second.append(["x", "y"])
    p = tmp_path / "sheet.xlsx"
    wb.save(str(p))

    out = extract_text(p)
    assert "apple" in out and "5" in out
    assert "--- sheet: Data ---" in out
    assert "--- sheet: More ---" in out


def test_pptx_extracts_text_and_notes(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Slide text here"
    slide.notes_slide.notes_text_frame.text = "speaker note"
    p = tmp_path / "deck.pptx"
    prs.save(str(p))

    out = extract_text(p)
    assert "Slide text here" in out
    assert "--- slide 1 ---" in out
    assert "speaker note" in out


def test_pdf_extracts_text_with_page_markers(tmp_path):
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=12)
    pdf.cell(0, 10, "Hello PDF content")
    p = tmp_path / "doc.pdf"
    pdf.output(str(p))

    out = extract_text(p)
    assert "Hello PDF content" in out
    assert "--- page 1 ---" in out


# ---- 加固回归(对抗审查) ----


def test_input_size_limit_rejects_large_file(tmp_path, monkeypatch):
    monkeypatch.setattr(extract_mod, "_MAX_INPUT_BYTES", 10)
    p = tmp_path / "big.txt"
    p.write_text("x" * 100)
    with pytest.raises(RuntimeError, match="too large"):
        extract_text(p)


def test_zip_bomb_uncompressed_limit(tmp_path, monkeypatch):
    import docx

    monkeypatch.setattr(extract_mod, "_MAX_UNCOMPRESSED_BYTES", 100)
    document = docx.Document()
    document.add_paragraph("content that uncompresses well beyond one hundred bytes")
    p = tmp_path / "doc.docx"
    document.save(str(p))
    with pytest.raises(RuntimeError, match="uncompressed"):
        extract_text(p)


def test_corrupt_docx_friendly_error_without_host_path(tmp_path):
    p = tmp_path / "broken.docx"
    p.write_bytes(b"this is plainly not a zip archive")
    with pytest.raises(RuntimeError) as exc:
        extract_text(p)
    msg = str(exc.value)
    assert "broken.docx" in msg
    assert str(tmp_path) not in msg  # 不把 host 绝对路径泄露给模型


def test_scanned_pdf_returns_blank_hint(tmp_path):
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()  # 空白页,无文字层(模拟扫描件/纯图)
    p = tmp_path / "scan.pdf"
    pdf.output(str(p))
    out = extract_text(p)
    assert "no extractable text" in out


def test_empty_xlsx_returns_blank_hint(tmp_path):
    from openpyxl import Workbook

    p = tmp_path / "empty.xlsx"
    Workbook().save(str(p))
    out = extract_text(p)
    assert "no extractable text" in out


def test_xxe_external_entity_not_resolved(tmp_path):
    """恶意 docx 注入外部实体 → 不应读到本地 secret(锁住三库 resolve_entities=False 不变量)。"""
    import docx

    secret = tmp_path / "secret.txt"
    secret.write_text("TOPSECRET_XXE_CANARY")

    base = tmp_path / "base.docx"
    docx.Document().save(str(base))
    malicious_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        f'<!DOCTYPE w:document [<!ENTITY xxe SYSTEM "file://{secret}">]>\n'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body><w:p><w:r><w:t>BEFORE-&xxe;-AFTER</w:t></w:r></w:p></w:body></w:document>"
    )
    evil = tmp_path / "evil.docx"
    with zipfile.ZipFile(base) as zin, zipfile.ZipFile(evil, "w") as zout:
        for item in zin.namelist():
            data = malicious_xml.encode() if item == "word/document.xml" else zin.read(item)
            zout.writestr(item, data)

    try:
        out = extract_text(evil)
    except RuntimeError:
        return  # 直接拒绝解析也可接受 —— 同样没有泄露
    assert "TOPSECRET_XXE_CANARY" not in out
